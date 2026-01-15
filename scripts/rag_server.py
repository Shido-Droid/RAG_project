import os
import sys
import json
import sqlite3
import datetime
from typing import List, Any
from fastapi import FastAPI, HTTPException, UploadFile, File
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse, StreamingResponse
from pydantic import BaseModel
import uvicorn
import io
from pypdf import PdfReader
from docx import Document
from pptx import Presentation
import pytesseract
from pdf2image import convert_from_bytes
from PIL import Image

from src.rag_app import (
    process_question, 
    add_document_to_kb, 
    clear_knowledge_base, 
    get_all_documents, 
    delete_document_from_kb, 
    document_exists, 
    analyze_document_content, 
    update_document_title, 
    explain_term
)

app = FastAPI(
    title="RAG Qwen Ultimate API",
    description="Qwen-based RAG System API with Intent Detection and Web Search",
    version="2.0"
)

# Reactフロントエンド(http://localhost:5173)からのアクセスを許可
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],  # 開発用: すべてのオリジンを許可
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# 静的ファイルのパス設定 (Project Root relative to scripts/rag_server.py)
BASE_DIR = os.path.dirname(os.path.abspath(__file__))
PROJECT_ROOT = os.path.dirname(BASE_DIR)
STATIC_DIR = os.path.join(PROJECT_ROOT, "static")
TEMPLATES_DIR = os.path.join(PROJECT_ROOT, "templates")

# staticディレクトリがない場合は作成（念のため）
os.makedirs(STATIC_DIR, exist_ok=True)

# DB Init
DB_PATH = os.path.join(BASE_DIR, "chat_history.db")

def init_db():
    with sqlite3.connect(DB_PATH) as conn:
        conn.execute('''CREATE TABLE IF NOT EXISTS history
                        (id INTEGER PRIMARY KEY AUTOINCREMENT,
                         sender TEXT,
                         text TEXT,
                         timestamp TEXT,
                         sources TEXT)''')

init_db()

# Static Files Mounting
from fastapi.staticfiles import StaticFiles
app.mount("/static", StaticFiles(directory=STATIC_DIR), name="static")

# --- Request/Response Models ---
class QueryRequest(BaseModel):
    question: str
    difficulty: str = "normal"

class DeleteRequest(BaseModel):
    filename: str

class UpdateTitleRequest(BaseModel):
    filename: str
    new_title: str

class ExplainRequest(BaseModel):
    term: str

class Source(BaseModel):
    title: str
    url: str

class QueryResponse(BaseModel):
    answer: str
    sources: List[Source]

# --- Endpoints ---
@app.get("/")
async def read_root():
    # templates/index.html を返す
    return FileResponse(os.path.join(TEMPLATES_DIR, "index.html"))

def get_pdf_sections(reader: PdfReader) -> List[tuple]:
    """PDFのアウトライン（しおり）からセクション情報を抽出する"""
    sections = []
    try:
        def _process_outline(outline):
            for item in outline:
                if isinstance(item, list):
                    _process_outline(item)
                elif hasattr(item, "page") and item.page is not None:
                    try:
                        # item.page は IndirectObject の場合があるためページ番号を取得
                        page_num = reader.get_page_number(item.page)
                        if page_num is not None and page_num >= 0:
                            sections.append((page_num, item.title))
                    except Exception:
                        pass
        
        if hasattr(reader, "outline") and reader.outline:
            _process_outline(reader.outline)
            
        # ページ順にソート
        sections.sort(key=lambda x: x[0])
    except Exception:
        pass # アウトラインがない、または読み取れない場合は無視
    return sections

@app.get("/api/history")
def get_history_endpoint():
    """チャット履歴を取得"""
    try:
        with sqlite3.connect(DB_PATH) as conn:
            conn.row_factory = sqlite3.Row
            cursor = conn.execute("SELECT sender, text, timestamp, sources FROM history ORDER BY id ASC")
            rows = cursor.fetchall()
            history = []
            for row in rows:
                msg = {
                    "sender": row["sender"],
                    "text": row["text"],
                    "timestamp": row["timestamp"]
                }
                if row["sources"]:
                    try:
                        msg["sources"] = json.loads(row["sources"])
                    except:
                        msg["sources"] = []
                history.append(msg)
        return history
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.delete("/api/history")
def clear_history_endpoint():
    """チャット履歴を削除"""
    try:
        with sqlite3.connect(DB_PATH) as conn:
            conn.execute("DELETE FROM history")
        return {"status": "success"}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/api/explain")
async def explain_endpoint(request: ExplainRequest):
    """専門用語の解説を生成して返します"""
    explanation = explain_term(request.term)
    return {"term": request.term, "explanation": explanation}

@app.post("/api/ask")
async def ask_endpoint(request: QueryRequest):
    """
    RAGシステムに質問を投げ、回答と参照元をストリーミング形式(NDJSON)で返します。
    """
    timestamp = datetime.datetime.now().strftime("%H:%M")
    
    # 1. 過去の履歴を取得 (直近3往復分程度)
    history = []
    try:
        with sqlite3.connect(DB_PATH) as conn:
            conn.row_factory = sqlite3.Row
            # 最新の6件を取得
            cursor = conn.execute("SELECT sender, text FROM history ORDER BY id DESC LIMIT 6")
            rows = cursor.fetchall()
            # 時系列順に戻す
            for row in reversed(rows):
                role = "user" if row["sender"] == "user" else "assistant"
                history.append({"role": role, "content": row["text"]})
    except Exception as e:
        print(f"DB Error (History fetch): {e}")

    # ユーザーの質問を保存
    try:
        with sqlite3.connect(DB_PATH) as conn:
            conn.execute("INSERT INTO history (sender, text, timestamp) VALUES (?, ?, ?)", 
                         ("user", request.question, timestamp))
    except Exception as e:
        print(f"DB Error (User): {e}")

    def generate():
        try:
            # 【改善】処理開始のステータスを即座に通知
            yield json.dumps({"type": "status", "content": "質問を分析し、ドキュメントを検索中..."}, ensure_ascii=False) + "\n"

            # 処理実行 (同期関数)
            result = process_question(request.question, history=history, difficulty=request.difficulty)
            
            # AIの回答を保存
            try:
                with sqlite3.connect(DB_PATH) as conn:
                    sources_json = json.dumps(result["sources"], ensure_ascii=False)
                    conn.execute("INSERT INTO history (sender, text, timestamp, sources) VALUES (?, ?, ?, ?)", 
                                 ("bot", result["answer"], timestamp, sources_json))
            except Exception as e:
                print(f"DB Error (Bot): {e}")

            # 【改善】生成完了ステータス
            yield json.dumps({"type": "status", "content": "回答を表示します"}, ensure_ascii=False) + "\n"

            # App.jsx が期待する形式で分割して送信
            yield json.dumps({"type": "sources", "content": result["sources"]}, ensure_ascii=False) + "\n"
            yield json.dumps({"type": "answer", "content": result["answer"]}, ensure_ascii=False) + "\n"
            yield json.dumps({"type": "done", "content": "completed"}, ensure_ascii=False) + "\n"
        except Exception as e:
            yield json.dumps({"type": "error", "content": f"エラーが発生しました: {str(e)}"}, ensure_ascii=False) + "\n"

    return StreamingResponse(generate(), media_type="application/x-ndjson")

@app.get("/api/documents")
async def get_documents_endpoint():
    """登録済みドキュメントの一覧を取得"""
    return {"documents": get_all_documents()}

@app.post("/api/reset")
async def reset_db_endpoint():
    """ChromaDBを初期化します"""
    try:
        clear_knowledge_base()
        return {"status": "success", "message": "Knowledge base cleared."}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/api/delete_document")
async def delete_document_endpoint(request: DeleteRequest):
    """指定されたドキュメントを削除します"""
    try:
        success = delete_document_from_kb(request.filename)
        if success:
            return {"status": "success", "message": f"Deleted {request.filename}"}
        else:
            raise HTTPException(status_code=500, detail="Failed to delete document")
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/api/update_title")
async def update_title_endpoint(request: UpdateTitleRequest):
    """ドキュメントのタイトルを更新します"""
    try:
        success = update_document_title(request.filename, request.new_title)
        if success:
            return {"status": "success", "message": f"Updated title for {request.filename}"}
        else:
            raise HTTPException(status_code=404, detail="Document not found or failed to update")
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/api/upload")
async def upload_file_endpoint(file: UploadFile = File(...), autorename: bool = True):
    """PDF/Word/PPT/Text/ImageファイルをアップロードしてRAGに取り込みます"""
    filename = file.filename
    if not filename:
        raise HTTPException(status_code=400, detail="Filename is required")

    # 重複チェック
    if document_exists(filename):
        if not autorename:
            raise HTTPException(status_code=409, detail=f"File '{filename}' already exists.")
        
        # 自動リネーム: filename_v2.ext, filename_v3.ext ...
        base, ext = os.path.splitext(filename)
        counter = 2
        while True:
            new_filename = f"{base}_v{counter}{ext}"
            if not document_exists(new_filename):
                filename = new_filename
                break
            counter += 1

    content = await file.read()
    file_obj = io.BytesIO(content)
    text = ""

    try:
        if filename.lower().endswith(".pdf"):
            reader = PdfReader(file_obj)
            sections = get_pdf_sections(reader)
            
            for i, page in enumerate(reader.pages):
                page_text = page.extract_text()
                if page_text:
                    # 現在のページに対応するセクションを探す
                    current_section = None
                    for start_page, title in sections:
                        if start_page <= i:
                            current_section = title
                        else:
                            break
                    
                    header = f"[Page {i+1}]"
                    if current_section:
                        header = f"[Section: {current_section}] {header}"
                    
                    text += f"{header} {page_text}\n"
            
            # テキストが抽出できなかった場合、OCRを試行 (スキャンPDF対策)
            if not text.strip():
                images = convert_from_bytes(content)
                for i, img in enumerate(images):
                    text += f"[Page {i+1}] {pytesseract.image_to_string(img, lang='jpn+eng')}\n"
        
        elif filename.lower().endswith((".jpg", ".jpeg", ".png", ".bmp", ".tiff")):
            image = Image.open(file_obj)
            text = pytesseract.image_to_string(image, lang="jpn+eng")
        
        elif filename.lower().endswith(".docx"):
            doc = Document(file_obj)
            for para in doc.paragraphs:
                text += para.text + "\n"
        
        elif filename.lower().endswith(".pptx"):
            prs = Presentation(file_obj)
            for slide in prs.slides:
                for shape in slide.shapes:
                    s: Any = shape
                    if hasattr(s, "text"):
                        text += s.text + "\n"
        
        elif filename.lower().endswith(".txt") or filename.lower().endswith(".md"):
            text = content.decode("utf-8", errors="ignore")
        
        else:
            return {"status": "error", "message": "Unsupported file type"}

        if not text.strip():
            msg = "No text extracted."
            if filename.lower().endswith(".pdf"):
                msg += " OCR was attempted but failed to extract text."
            return {"status": "error", "message": msg}

        # ドキュメント分析 (要約・タイトル・キーワード)
        doc_meta = analyze_document_content(text)

        # DBに追加
        add_document_to_kb(text, source=filename, doc_metadata=doc_meta)
        return {"status": "success", "message": f"Imported {filename} ({len(text)} chars)"}

    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Failed to process file: {str(e)}")

if __name__ == "__main__":
    # 開発用サーバー起動
    print(f"Starting RAG Server on http://0.0.0.0:8000 (Docs: http://localhost:8000/docs)")
    uvicorn.run(app, host="0.0.0.0", port=8000)